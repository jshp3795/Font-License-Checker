from pptx import Presentation

prs = Presentation("a.pptx")

fonts = []

for slide in prs.slides:
    for shape in slide.shapes:
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                fonts.append(run.font.name)

print(fonts)
